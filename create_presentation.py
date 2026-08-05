from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ============ SLIDE 1: TITLE ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0A, 0x1A, 0x2E)

    box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "MediCare AI - Healthcare Chatbot"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    box2 = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1))
    p2 = box2.text_frame.paragraphs[0]
    p2.text = "AI-Powered Healthcare Assistant with RAG and Conversation Memory"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0x48, 0xCA, 0xE4)
    p2.alignment = PP_ALIGN.CENTER

    box3 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    p3 = box3.text_frame.paragraphs[0]
    p3.text = "OpenRouter LLM | LangChain | FAISS | Streamlit"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
    p3.alignment = PP_ALIGN.CENTER

    box4 = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(0.7))
    p4 = box4.text_frame.paragraphs[0]
    p4.text = "AI Engineer Technical Assignment | 2026"
    p4.font.size = Pt(14)
    p4.font.color.rgb = RGBColor(0x90, 0xA4, 0xAE)
    p4.alignment = PP_ALIGN.CENTER

    # ============ SLIDE 2: SYSTEM ARCHITECTURE ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0A, 0x1A, 0x2E)

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = "System Architecture"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    layers = [
        ("STREAMLIT FRONTEND", "Chat UI | Sidebar Controls | Quick Questions | Disclaimer Banner", Inches(1.2), RGBColor(0x1B, 0x5E, 0x8A)),
        ("HEALTHCARE CHATBOT CORE", "Query Classifier | Response Generator | Guardrails | Disclaimers", Inches(2.5), RGBColor(0x1B, 0x3A, 0x4B)),
        ("FAISS VECTOR STORE", "20 Medical Docs | all-MiniLM-L6-v2 Embeddings | Similarity Search", Inches(4.0), RGBColor(0x3D, 0x1F, 0x5C)),
        ("OPENROUTER LLM (GPT-OSS-20B)", "Temperature 0.4 | Max Tokens 2048 | OpenAI-Compatible API", Inches(4.0), RGBColor(0x5C, 0x1F, 0x1F)),
        ("CONVERSATION MEMORY", "Sliding Window (10 exchanges) | 20 Message History | Context-Aware", Inches(4.0), RGBColor(0x1F, 0x5C, 0x1F)),
    ]

    # Top two layers full width
    for title, desc, y, color in layers[:2]:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(9), Inches(0.9), )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        p2.alignment = PP_ALIGN.CENTER

    # Bottom three in columns
    for i, (title, desc, y, color) in enumerate(layers[2:]):
        x = Inches(0.5 + i * 3.15)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(9)
        p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        p2.alignment = PP_ALIGN.CENTER

    box5 = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.5))
    p5 = box5.text_frame.paragraphs[0]
    p5.text = "Flow: User Query -> Classification -> RAG Retrieval + LLM -> Safety Checks -> Response"
    p5.font.size = Pt(11)
    p5.font.color.rgb = RGBColor(0x48, 0xCA, 0xE4)
    p5.alignment = PP_ALIGN.CENTER

    # ============ SLIDE 3: TECH STACK & LLM ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0A, 0x1A, 0x2E)

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = "Tech Stack and LLM Selection"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    table_data = [
        ["Component", "Technology", "Rationale"],
        ["Frontend", "Streamlit", "Rapid prototyping, built-in chat, Python-native"],
        ["LLM", "OpenRouter (GPT-OSS-20B)", "Free tier, fast, OpenAI-compatible API"],
        ["Orchestration", "LangChain", "Mature framework for chains and RAG"],
        ["Vector Store", "FAISS", "Efficient CPU-based similarity search"],
        ["Embeddings", "all-MiniLM-L6-v2", "Lightweight, fast, good semantics"],
        ["Language", "Python 3.9+", "Best AI/ML ecosystem"],
    ]

    y = Inches(1.3)
    for i, row in enumerate(table_data):
        is_header = i == 0
        color = RGBColor(0x00, 0x77, 0xB6) if is_header else (RGBColor(0x12, 0x1E, 0x30) if i % 2 == 1 else RGBColor(0x0F, 0x18, 0x28))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(9), Inches(0.55))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].text = f"  {row[0]:<25} {row[1]:<30} {row[2]}"
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = is_header
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        y += Inches(0.6)

    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.8))
    p2 = box2.text_frame.paragraphs[0]
    p2.text = "LLM Rationale: GPT-OSS-20B via OpenRouter offers free access, fast inference, and OpenAI-compatible API for easy integration with LangChain."
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)

    # ============ SLIDE 4: PROMPT ENGINEERING ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0A, 0x1A, 0x2E)

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = "Prompt Engineering and Additional Components"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Left column
    left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(4.3), Inches(5))
    left.fill.solid()
    left.fill.fore_color.rgb = RGBColor(0x12, 0x1E, 0x30)
    left.line.fill.background()
    ltf = left.text_frame
    ltf.word_wrap = True
    ltf.paragraphs[0].text = "Prompt Strategy"
    ltf.paragraphs[0].font.size = Pt(16)
    ltf.paragraphs[0].font.bold = True
    ltf.paragraphs[0].font.color.rgb = RGBColor(0x48, 0xCA, 0xE4)

    strategies = [
        "Role Prompting - Empathetic healthcare assistant",
        "Constraint Setting - No diagnoses or prescriptions",
        "Context Injection - RAG docs in system prompt",
        "Guardrail Prompting - Safety in every prompt",
        "Response Templates - Emergency/greeting/OOS",
        "Chain-of-Thought - Classification then response",
        "Medical Disclaimers - Auto-appended to responses",
    ]
    for s in strategies:
        p = ltf.add_paragraph()
        p.text = f"  {s}"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        p.space_before = Pt(6)

    # Right column
    right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.3), Inches(4.3), Inches(5))
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(0x12, 0x1E, 0x30)
    right.line.fill.background()
    rtf = right.text_frame
    rtf.word_wrap = True
    rtf.paragraphs[0].text = "Additional Components"
    rtf.paragraphs[0].font.size = Pt(16)
    rtf.paragraphs[0].font.bold = True
    rtf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xA5, 0x00)

    components = [
        ("RAG Pipeline", "FAISS + HuggingFace for context retrieval"),
        ("Vector Database", "FAISS with 20 curated medical docs"),
        ("Conversation Memory", "Sliding window buffer + history"),
        ("Emergency Detection", "Keyword-based safety response"),
        ("Scope Enforcement", "Redirects non-health queries"),
        ("Source Citations", "References to WHO, CDC, etc."),
        ("Query Classification", "Routes queries to handlers"),
    ]
    for title, desc in components:
        p = rtf.add_paragraph()
        p.text = f"  {title}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xDD, 0x77)
        p.space_before = Pt(6)
        p2 = rtf.add_paragraph()
        p2.text = f"    {desc}"
        p2.font.size = Pt(9)
        p2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)

    # ============ SLIDE 5: CHALLENGES ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0A, 0x1A, 0x2E)

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = "Challenges, Solutions and Workflow"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    challenges = [
        ("LLM Hallucination", "RAG grounds responses in verified medical knowledge base"),
        ("Medical Safety", "Multi-layer guardrails: emergency detection, diagnosis prevention, disclaimers"),
        ("Context in Follow-ups", "ConversationBufferWindowMemory maintains last 10 exchanges"),
        ("Response Relevance", "Query classification routes to appropriate handlers"),
        ("Source Attribution", "Knowledge base includes source metadata for citations"),
    ]

    y = Inches(1.3)
    for title, solution in challenges:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(9), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x12, 0x1E, 0x30)
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = f"Challenge: {title}"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0x66, 0x66)
        p2 = tf.add_paragraph()
        p2.text = f"Solution: {solution}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(0x77, 0xDD, 0x77)
        p2.space_before = Pt(4)
        y += Inches(1.0)

    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    p2 = box2.text_frame.paragraphs[0]
    p2.text = "Workflow: User Input -> Classification -> RAG + LLM -> Safety Validation -> Response with Citations"
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0x48, 0xCA, 0xE4)
    p2.alignment = PP_ALIGN.CENTER

    # Save
    os.makedirs('docs', exist_ok=True)
    prs.save('docs/architecture_presentation.pptx')
    print('Presentation saved: docs/architecture_presentation.pptx')

if __name__ == '__main__':
    create_presentation()
